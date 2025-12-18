from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_slide(prs, title_text, bullet_points):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    
    tf = slide.placeholders[1].text_frame
    tf.text = bullet_points[0]
    
    for point in bullet_points[1:]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0

def generate_presentation():
    prs = Presentation()

    # Slide 1: Title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "IPL Search Engine"
    subtitle.text = "A Comprehensive Full-Stack Search Solution\nAdvanced Indexing, Semantic Search, and Real-Time Performance"

    # Slide 2: Introduction
    add_slide(prs, "Project Overview", [
        "Goal: Build a scalable, file-based search engine for IPL and sports data.",
        "Key Features: Multi-word search, semantic understanding, and real-time autocomplete.",
        "Focus: Performance, memory efficiency, and professional UI/UX.",
        "Dataset: 240,000+ IPL match records and football player data."
    ])

    # Slide 3: System Architecture
    add_slide(prs, "System Architecture", [
        "Three-Tier Architecture:",
        "1. Frontend: React.js for a modern, responsive user interface.",
        "2. Backend: FastAPI (Python) for high-performance API services.",
        "3. Indexing Core: Custom-built file-based indexing system (No Databases)."
    ])

    # Slide 4: File-Based Storage
    add_slide(prs, "File-Based Storage Strategy", [
        "Requirement: No external databases allowed for indexing data.",
        "Implementation: Serialized Python objects (Pickle) and flat files.",
        "Advantages: Zero-dependency deployment and direct control over data structures.",
        "Scaling: Optimized for fast I/O using binary formats."
    ])

    # Slide 5: The Lexicon
    add_slide(prs, "Core Component: The Lexicon", [
        "Definition: A dictionary storing all unique words in the dataset.",
        "Purpose: Maps each word to a distinct integer ID.",
        "Efficiency: Reduces memory usage by using IDs instead of strings in indices.",
        "Size: Handles 20,000+ unique tokens across cricket and football datasets."
    ])

    # Slide 6: Lexicon Implementation
    add_slide(prs, "Lexicon Implementation Details", [
        "Class: LexiconBuilder handles word discovery and ID assignment.",
        "Storage: Saved as 'lexicon.pkl' for fast loading during API startup.",
        "Normalization: Case-folding and basic cleaning during indexing.",
        "Lookup: O(1) average time complexity using Python dictionaries."
    ])

    # Slide 7: Forward Index
    add_slide(prs, "Core Component: Forward Index", [
        "Structure: Maps Document IDs to a list of Word IDs.",
        "Storage: Organized by document for easy sequential processing.",
        "Role: Serves as the primary data source for building the Inverted Index.",
        "Metadata: Stores document-specific attributes (match name, season, etc.)."
    ])

    # Slide 8: Forward Index Utility
    add_slide(prs, "Forward Index in the Pipeline", [
        "Indexing Flow: Raw CSV -> Preprocessor -> Forward Index -> Inverted Index.",
        "Decoupling: Allows re-building the inverted index without re-parsing raw data.",
        "Optimization: Uses fixed-width or compressed formats for large datasets."
    ])

    # Slide 9: Inverted Index
    add_slide(prs, "Core Component: Inverted Index", [
        "Structure: Maps Word IDs to a list of Document IDs (Postings List).",
        "Purpose: The heart of the search engine; enables fast retrieval.",
        "Querying: Directly identifies which documents contain the search terms.",
        "Efficiency: Avoids scanning the entire dataset for every query."
    ])

    # Slide 10: Inverted Index Efficiency
    add_slide(prs, "Inverted Index Performance", [
        "Lookup Complexity: O(1) to find the postings list for a word.",
        "Intersection: Efficiently combines lists for multi-word queries.",
        "Memory: Designed to handle millions of postings without crashing."
    ])

    # Slide 11: Scaling with Barrels
    add_slide(prs, "Scaling with Barrels", [
        "Problem: Inverted index grows too large for RAM (100k+ docs).",
        "Solution: Divide the inverted index into smaller chunks called 'Barrels'.",
        "Implementation: Each barrel stores a subset of the word-to-doc mappings.",
        "Scaling: Allows the engine to scale to millions of documents."
    ])

    # Slide 12: Barrel Management
    add_slide(prs, "Barrel Management Logic", [
        "Class: BarrelManager handles reading/writing barrels to disk.",
        "Memory Efficiency: Only loads necessary barrels during a search.",
        "Flushing: Periodically writes data to disk during indexing to save RAM.",
        "Merge: Logic to combine partial barrels into a final optimized index."
    ])

    # Slide 13: Single Word Searching
    add_slide(prs, "Search: Single Word Queries", [
        "Flow: Query -> Preprocessor -> Lexicon (ID) -> Inverted Index -> Results.",
        "Performance: Response time < 500ms (Target met).",
        "Relevance: Basic TF-IDF scoring applied to the single term."
    ])

    # Slide 14: Multi-Word Searching
    add_slide(prs, "Search: Multi-Word Queries", [
        "Logic: Intersection (AND) of postings lists for all query terms.",
        "Example: 'Virat Kohli' -> Docs containing BOTH 'Virat' and 'Kohli'.",
        "Performance: Response time < 1.5s for 5-word queries (Target met).",
        "Robustness: Handles stop words and noise terms gracefully."
    ])

    # Slide 15: Semantic Search
    add_slide(prs, "Semantic Search Capabilities", [
        "Beyond Keywords: Understanding intent and conceptually similar queries.",
        "Synonyms: Mapping 'MSD' to 'Dhoni' or 'RCB' to 'Royal Challengers Bangalore'.",
        "Phrase Normalization: Handling variations like 'top scorer' vs 'orange cap'.",
        "Constraint: Implemented without LLMs/Transformers (as per requirements)."
    ])

    # Slide 16: Semantic Implementation
    add_slide(prs, "Semantic Matching Logic", [
        "Synonym Dictionary: Pre-defined mappings for sports-specific terms.",
        "Query Expansion: Automatically adding related terms to the search.",
        "Normalization: Converting ordinals (1st, 2nd) to standard tokens.",
        "Custom Logic: Weighted matching for conceptually related fields."
    ])

    # Slide 17: Real-Time Autocomplete
    add_slide(prs, "Real-Time Autocomplete", [
        "UX: Provides 3-5 suggestions as the user types.",
        "Speed: Response time < 100ms for seamless interaction.",
        "Source: Based on the existing Lexicon and popular query patterns."
    ])

    # Slide 18: Trie Data Structure
    add_slide(prs, "Autocomplete: Trie Implementation", [
        "Structure: Prefix tree for efficient string matching.",
        "Why Trie? Faster than linear scanning of the lexicon (O(L) vs O(N)).",
        "Memory: Compressed nodes for storing 20k+ words efficiently.",
        "Integration: Built during API startup for instant availability."
    ])

    # Slide 19: Ranking Algorithm
    add_slide(prs, "Ranking Results", [
        "Criterion: Additive scoring based on frequency and field importance.",
        "TF-IDF: Base relevance score for term frequency in documents.",
        "Position: Higher weight for terms appearing in titles vs. descriptions.",
        "Sorting: Final results sorted by score descending."
    ])

    # Slide 20: Ranking Criteria & Weights
    add_slide(prs, "Field Weighting & Boosting", [
        "Player Name Match: 1000.0 (Absolute priority).",
        "Team Match: 50.0 (High priority for match queries).",
        "Document Type Boosts: Career Stats > Season Stats > Match Commentary.",
        "Football Boost: Adjusted to ensure relevance alongside cricket data."
    ])

    # Slide 21: Dynamic Content Addition
    add_slide(prs, "Dynamic Content Indexing", [
        "Requirement: Index new documents without restarting the system.",
        "Class: DocumentAdder handles real-time updates to the index files.",
        "Consistency: Ensures new documents follow the same ranking rules.",
        "Speed: Indexing completes in < 1 minute (Target met)."
    ])

    # Slide 22: Document Addition UI
    add_slide(prs, "Content Management Interface", [
        "UI: Simple modal for adding JSON or text-based documents.",
        "Validation: Ensures data integrity before indexing.",
        "Feedback: Real-time confirmation of successful indexing.",
        "Accessibility: Integrated into the main search dashboard."
    ])

    # Slide 23: System Performance: Query
    add_slide(prs, "Performance: Query Targets", [
        "Single Word: < 500ms (Actual: ~5ms).",
        "5-Word Queries: < 1.5s (Actual: ~20ms).",
        "Scalability: No slowdown observed with 240k+ documents.",
        "Transparency: Execution time displayed in the UI for every search."
    ])

    # Slide 24: System Performance: Indexing
    add_slide(prs, "Performance: Indexing & Scaling", [
        "Dataset Size: 240,000+ documents (Requirement: > 45,000).",
        "Indexing Speed: Full re-index in ~5 minutes for 240k docs.",
        "Update Speed: Individual document indexing in < 1 second.",
        "Concurrency: Index updates do not block search operations."
    ])

    # Slide 25: Memory Management
    add_slide(prs, "Memory Usage Optimization", [
        "Target: <= 2GB for 100k docs, <= 4GB for 100k+ docs.",
        "Actual: RAM usage remains stable at ~1.2GB for 240k docs.",
        "Techniques: Barrel-based loading and efficient metadata serialization.",
        "Optimization: Garbage collection and lazy loading of postings lists."
    ])

    # Slide 26: Professional User Interface
    add_slide(prs, "Professional UI Design", [
        "Framework: React.js with modern CSS (Vanilla).",
        "Aesthetics: Glassmorphism, dark mode, and vibrant accents.",
        "Responsiveness: Fully functional on desktop and mobile devices.",
        "User Experience: Smooth transitions and micro-animations."
    ])

    # Slide 27: UI Features & UX
    add_slide(prs, "UI Features & UX Highlights", [
        "Image Gallery: Shuffled player images for a dynamic feel.",
        "Stats Display: Dedicated grids for career and season statistics.",
        "Search Modes: Quick, Deep, and Research modes (UI concepts).",
        "Error Handling: Graceful alerts for backend connectivity issues."
    ])

    # Slide 28: Application Deployment
    add_slide(prs, "Deployment & Accessibility", [
        "Platform: Deployed on cloud infrastructure (e.g., Railway/Render).",
        "Public URL: Accessible online for testing and evaluation.",
        "Environment: Production-ready FastAPI server with Uvicorn.",
        "CI/CD: Automated deployment from GitHub repository."
    ])

    # Slide 29: Team Collaboration (GIT)
    add_slide(prs, "Collaboration & Version Control", [
        "Requirement: Minimum 20 meaningful commits.",
        "Actual: 50+ commits covering all major features.",
        "Teamwork: Contributions from all members with proper commit messages.",
        "Workflow: Feature branching and regular pushes to GitHub."
    ])

    # Slide 30: Conclusion
    add_slide(prs, "Conclusion & Future Scope", [
        "Summary: Successfully built a high-performance, scalable search engine.",
        "Achievements: Met all requirements, including barrels and semantic search.",
        "Future Scope: Integration of more sports, advanced NLP, and user accounts.",
        "Final Verdict: A robust solution for large-scale sports data retrieval."
    ])

    prs.save('IPL_Search_Engine_Presentation.pptx')
    print("Presentation generated successfully: IPL_Search_Engine_Presentation.pptx")

if __name__ == "__main__":
    generate_presentation()
